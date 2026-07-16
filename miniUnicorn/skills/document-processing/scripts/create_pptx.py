#!/usr/bin/env python3
"""用 python-pptx 创建 PowerPoint 文档。

用法:
    # 基本演示
    python create_pptx.py output.pptx --title "季度报告" --slides '[{"title":"Q1","content":"营收增长15%"},{"title":"Q2","content":"拓展新市场"}]'

    # 带项目符号
    python create_pptx.py output.pptx --title "项目计划" --slides '[{"title":"目标","bullets":["增长20%","降低成本"]}]'
"""

from __future__ import annotations

import argparse
import json
import sys
from pathlib import Path


def create_pptx(
    output_path: Path,
    *,
    title: str | None = None,
    slides_json: str | None = None,
) -> None:
    """创建 .pptx 文档。"""
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
    except ImportError:
        print("错误:需要 python-pptx 库。安装:pip install python-pptx", file=sys.stderr)
        sys.exit(1)

    prs = Presentation()

    # 默认 16:9 尺寸
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # 标题页
    if title:
        slide_layout = prs.slide_layouts[0]  # Title Slide
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = title
        if len(slide.placeholders) > 1:
            slide.placeholders[1].text = ""

    # 内容页
    if slides_json:
        try:
            slides_data = json.loads(slides_json)
        except json.JSONDecodeError as e:
            print(f"错误:--slides JSON 解析失败: {e}", file=sys.stderr)
            sys.exit(1)

        for slide_data in slides_data:
            slide_layout = prs.slide_layouts[1]  # Title and Content
            slide = prs.slides.add_slide(slide_layout)
            if slide_data.get("title"):
                slide.shapes.title.text = slide_data["title"]
            if slide_data.get("content"):
                if len(slide.placeholders) > 1:
                    slide.placeholders[1].text = slide_data["content"]
            if slide_data.get("bullets"):
                if len(slide.placeholders) > 1:
                    tf = slide.placeholders[1].text_frame
                    tf.clear()
                    for i, bullet in enumerate(slide_data["bullets"]):
                        if i == 0:
                            tf.paragraphs[0].text = bullet
                        else:
                            p = tf.add_paragraph()
                            p.text = bullet

    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))
    print(f"已创建: {output_path} ({len(prs.slides)} 页)")


def main() -> None:
    parser = argparse.ArgumentParser(description="创建 PowerPoint 文档")
    parser.add_argument("output", help="输出 .pptx 路径")
    parser.add_argument("--title", help="演示文稿标题")
    parser.add_argument(
        "--slides",
        help='幻灯片 JSON,如 \'[{"title":"页1","content":"内容","bullets":["a","b"]}]\'',
    )
    args = parser.parse_args()

    create_pptx(
        Path(args.output).resolve(),
        title=args.title,
        slides_json=args.slides,
    )


if __name__ == "__main__":
    main()
