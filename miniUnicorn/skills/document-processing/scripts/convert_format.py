#!/usr/bin/env python3
"""文档格式转换(纯 Python,无外部 CLI 依赖)。

用法:
    python convert_format.py input.docx --to txt
    python convert_format.py input.docx --to pdf
    python convert_format.py input.xlsx --to csv
    python convert_format.py input.pptx --to txt

支持:
    .docx → .txt/.md/.pdf(基础)
    .xlsx → .csv(每个 sheet 一个文件)/.txt
    .pptx → .txt/.md
    .pdf  → .txt

注:.docx→.pdf 的纯 Python 方案仅支持基础排版。
"""

from __future__ import annotations

import argparse
import csv
import sys
from pathlib import Path


def docx_to_text(path: Path) -> str:
    """提取 .docx 全部文本。"""
    from docx import Document
    doc = Document(path)
    parts: list[str] = []
    for para in doc.paragraphs:
        if para.text.strip():
            style = para.style.name if para.style else ""
            if "Heading 1" in style:
                parts.append(f"# {para.text}")
            elif "Heading 2" in style:
                parts.append(f"## {para.text}")
            elif "Heading 3" in style:
                parts.append(f"### {para.text}")
            else:
                parts.append(para.text)
    # 表格
    for table in doc.tables:
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells]
            parts.append(" | ".join(cells))
    return "\n\n".join(parts)


def xlsx_to_csv(path: Path, output_dir: Path, stem: str) -> list[Path]:
    """转换 .xlsx 为多个 CSV(每个 sheet 一个)。"""
    from openpyxl import load_workbook
    wb = load_workbook(path, read_only=True, data_only=True)
    outputs: list[Path] = []
    try:
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            out_file = output_dir / f"{stem}_{sheet_name}.csv"
            with open(out_file, "w", newline="", encoding="utf-8") as f:
                writer = csv.writer(f)
                for row in ws.iter_rows(values_only=True):
                    writer.writerow([str(c) if c is not None else "" for c in row])
            outputs.append(out_file)
    finally:
        wb.close()
    return outputs


def pptx_to_text(path: Path) -> str:
    """提取 .pptx 全部文本。"""
    from pptx import Presentation
    prs = Presentation(path)
    parts: list[str] = []
    for i, slide in enumerate(prs.slides, 1):
        slide_texts: list[str] = [f"--- Slide {i} ---"]
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if text:
                    slide_texts.append(text)
            if shape.has_table:
                for row in shape.table.rows:
                    cells = [cell.text.strip() for cell in row.cells]
                    slide_texts.append(" | ".join(cells))
        if len(slide_texts) > 1:
            parts.append("\n".join(slide_texts))
    return "\n\n".join(parts)


def pdf_to_text(path: Path) -> str:
    """提取 PDF 全部文本。"""
    from pypdf import PdfReader
    reader = PdfReader(path)
    parts: list[str] = []
    for i, page in enumerate(reader.pages, 1):
        text = page.extract_text() or ""
        parts.append(f"--- Page {i} ---\n{text}")
    return "\n\n".join(parts)


def docx_to_pdf(path: Path, output: Path) -> None:
    """将 .docx 转为 .pdf(基础版,用 reportlab + python-docx)。

    注意:这是纯 Python 方案,仅支持文字 + 段落。复杂排版(字体/图片/表格样式)
    建议用户用 LibreOffice(soffice --convert-to pdf)手动转换。
    """
    try:
        from reportlab.lib.pagesizes import A4
        from reportlab.lib.styles import ParagraphStyle, getSampleStyleSheet
        from reportlab.lib.units import cm
        from reportlab.platypus import PageBreak, Paragraph, SimpleDocTemplate, Spacer
    except ImportError:
        print("错误:需要 reportlab 库。安装:pip install reportlab", file=sys.stderr)
        sys.exit(1)
    from docx import Document

    doc = Document(path)
    styles = getSampleStyleSheet()
    body_style = ParagraphStyle(
        "Body",
        parent=styles["Normal"],
        fontName="Helvetica",
        fontSize=11,
        leading=16,
        spaceAfter=8,
    )
    heading_style = ParagraphStyle(
        "Heading",
        parent=styles["Heading1"],
        fontSize=16,
        spaceAfter=12,
    )

    story = []
    for para in doc.paragraphs:
        if not para.text.strip():
            continue
        style_name = para.style.name if para.style else ""
        if "Heading" in style_name:
            story.append(Paragraph(para.text, heading_style))
        else:
            # 转义 XML 特殊字符
            text = para.text.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")
            story.append(Paragraph(text, body_style))
        story.append(Spacer(1, 4))

    output.parent.mkdir(parents=True, exist_ok=True)
    pdf_doc = SimpleDocTemplate(str(output), pagesize=A4, topMargin=2*cm, bottomMargin=2*cm)
    pdf_doc.build(story)


def convert(input_path: Path, target_format: str) -> None:
    """执行格式转换。"""
    ext = input_path.suffix.lower()
    stem = input_path.stem
    output_dir = input_path.parent

    if ext == ".docx":
        if target_format in ("txt", "md"):
            output = output_dir / f"{stem}.{target_format}"
            output.write_text(docx_to_text(input_path), encoding="utf-8")
            print(f"已转换: {input_path} → {output}")
        elif target_format == "pdf":
            output = output_dir / f"{stem}.pdf"
            docx_to_pdf(input_path, output)
            print(f"已转换: {input_path} → {output}")
            print("注:纯 Python 方案仅支持基础排版,复杂样式建议用 LibreOffice。")
        else:
            print(f"错误:不支持 .docx → .{target_format}", file=sys.stderr)
            sys.exit(1)

    elif ext == ".xlsx":
        if target_format == "csv":
            outputs = xlsx_to_csv(input_path, output_dir, stem)
            for o in outputs:
                print(f"已转换: {input_path} → {o}")
        elif target_format == "txt":
            from openpyxl import load_workbook
            wb = load_workbook(input_path, read_only=True, data_only=True)
            parts: list[str] = []
            try:
                for sheet_name in wb.sheetnames:
                    ws = wb[sheet_name]
                    parts.append(f"--- Sheet: {sheet_name} ---")
                    for row in ws.iter_rows(values_only=True):
                        parts.append("\t".join(str(c) if c is not None else "" for c in row))
            finally:
                wb.close()
            output = output_dir / f"{stem}.txt"
            output.write_text("\n".join(parts), encoding="utf-8")
            print(f"已转换: {input_path} → {output}")
        else:
            print(f"错误:不支持 .xlsx → .{target_format}", file=sys.stderr)
            sys.exit(1)

    elif ext == ".pptx":
        if target_format in ("txt", "md"):
            output = output_dir / f"{stem}.{target_format}"
            output.write_text(pptx_to_text(input_path), encoding="utf-8")
            print(f"已转换: {input_path} → {output}")
        else:
            print(f"错误:不支持 .pptx → .{target_format}", file=sys.stderr)
            sys.exit(1)

    elif ext == ".pdf":
        if target_format == "txt":
            output = output_dir / f"{stem}.txt"
            output.write_text(pdf_to_text(input_path), encoding="utf-8")
            print(f"已转换: {input_path} → {output}")
        else:
            print(f"错误:不支持 .pdf → .{target_format}", file=sys.stderr)
            sys.exit(1)

    else:
        print(f"错误:不支持的输入格式: {ext}", file=sys.stderr)
        sys.exit(1)


def main() -> None:
    parser = argparse.ArgumentParser(description="文档格式转换")
    parser.add_argument("input", help="输入文件路径")
    parser.add_argument("--to", required=True, choices=["txt", "md", "csv", "pdf"], help="目标格式")
    args = parser.parse_args()

    input_path = Path(args.input).resolve()
    if not input_path.exists():
        print(f"错误:文件不存在: {input_path}", file=sys.stderr)
        sys.exit(1)

    convert(input_path, args.to)


if __name__ == "__main__":
    main()
