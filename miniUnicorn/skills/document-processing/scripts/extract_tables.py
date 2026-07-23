#!/usr/bin/env python3
"""从文档中提取表格,输出为 JSON 或 CSV。

用法:
    python extract_tables.py input.docx --format json
    python extract_tables.py input.xlsx --format json
    python extract_tables.py input.pdf --format csv

JSON 输出格式:
    [
      {
        "sheet": "Sheet1",      # xlsx 的 sheet 名
        "page": 1,              # pdf 的页码(docx 无此字段)
        "headers": ["姓名", "年龄"],
        "rows": [["张三", 25], ["李四", 30]]
      }
    ]
"""

from __future__ import annotations

import argparse
import csv
import json
import sys
from pathlib import Path


def extract_docx_tables(path: Path) -> list[dict]:
    """从 .docx 提取所有表格。"""
    from docx import Document
    doc = Document(path)
    tables: list[dict] = []
    for i, table in enumerate(doc.tables, 1):
        rows = []
        for row in table.rows:
            rows.append([cell.text.strip() for cell in row.cells])
        if not rows:
            continue
        headers = rows[0] if rows else []
        tables.append({
            "table_index": i,
            "headers": headers,
            "rows": rows[1:] if len(rows) > 1 else [],
            "all_rows": rows,
        })
    return tables


def extract_xlsx_tables(path: Path) -> list[dict]:
    """从 .xlsx 提取所有 sheet 的表格数据。"""
    from openpyxl import load_workbook
    wb = load_workbook(path, read_only=True, data_only=True)
    tables: list[dict] = []
    try:
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            rows = []
            for row in ws.iter_rows(values_only=True):
                rows.append([str(c) if c is not None else "" for c in row])
            if not rows:
                continue
            # 过滤空行
            rows = [r for r in rows if any(c.strip() for c in r)]
            if not rows:
                continue
            headers = rows[0]
            tables.append({
                "sheet": sheet_name,
                "headers": headers,
                "rows": rows[1:] if len(rows) > 1 else [],
                "all_rows": rows,
            })
    finally:
        wb.close()
    return tables


def extract_pdf_tables(path: Path) -> list[dict]:
    """从 .pdf 提取表格(用 pdfplumber)。"""
    try:
        import pdfplumber
    except ImportError:
        print("错误:需要 pdfplumber 库。安装:pip install pdfplumber", file=sys.stderr)
        sys.exit(1)

    tables: list[dict] = []
    with pdfplumber.open(path) as pdf:
        for i, page in enumerate(pdf.pages, 1):
            page_tables = page.extract_tables() or []
            for j, table in enumerate(page_tables, 1):
                if not table:
                    continue
                rows = [[cell.strip() if cell else "" for cell in row] for row in table]
                headers = rows[0] if rows else []
                tables.append({
                    "page": i,
                    "table_index": j,
                    "headers": headers,
                    "rows": rows[1:] if len(rows) > 1 else [],
                    "all_rows": rows,
                })
    return tables


def extract_pptx_tables(path: Path) -> list[dict]:
    """从 .pptx 提取所有表格。"""
    from pptx import Presentation
    prs = Presentation(path)
    tables: list[dict] = []
    for i, slide in enumerate(prs.slides, 1):
        for shape in slide.shapes:
            if not shape.has_table:
                continue
            table = shape.table
            rows = []
            for row in table.rows:
                rows.append([cell.text.strip() for cell in row.cells])
            if not rows:
                continue
            headers = rows[0]
            tables.append({
                "slide": i,
                "headers": headers,
                "rows": rows[1:] if len(rows) > 1 else [],
                "all_rows": rows,
            })
    return tables


def output_json(tables: list[dict]) -> None:
    print(json.dumps(tables, ensure_ascii=False, indent=2))


def output_csv(tables: list[dict], input_stem: str) -> None:
    """每个表格输出一个 CSV 文件。"""
    if not tables:
        print("未找到表格。")
        return
    for i, table in enumerate(tables, 1):
        loc = table.get("sheet") or table.get("page") or table.get("slide") or i
        filename = f"{input_stem}_table_{loc}.csv"
        all_rows = [table["headers"]] + table["rows"]
        with open(filename, "w", newline="", encoding="utf-8") as f:
            writer = csv.writer(f)
            for row in all_rows:
                writer.writerow(row)
        print(f"已输出: {filename} ({len(table['rows'])} 行)")


def main() -> None:
    parser = argparse.ArgumentParser(description="从文档提取表格")
    parser.add_argument("input", help="输入文件")
    parser.add_argument("--format", choices=["json", "csv"], default="json", help="输出格式")
    args = parser.parse_args()

    input_path = Path(args.input).resolve()
    if not input_path.exists():
        print(f"错误:文件不存在: {input_path}", file=sys.stderr)
        sys.exit(1)

    ext = input_path.suffix.lower()

    if ext == ".docx":
        tables = extract_docx_tables(input_path)
    elif ext == ".xlsx":
        tables = extract_xlsx_tables(input_path)
    elif ext == ".pdf":
        tables = extract_pdf_tables(input_path)
    elif ext == ".pptx":
        tables = extract_pptx_tables(input_path)
    else:
        print(f"错误:不支持的格式: {ext}", file=sys.stderr)
        sys.exit(1)

    if not tables:
        print("未找到表格。")
        return

    if args.format == "json":
        output_json(tables)
    else:
        output_csv(tables, input_path.stem)


if __name__ == "__main__":
    main()
